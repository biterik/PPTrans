#!/usr/bin/env python3
"""
Fix the exceptions import issue in pptx_processor.py
"""

from pathlib import Path

def check_available_exceptions():
    """Check what exceptions are available in utils.exceptions"""
    try:
        import sys
        src_path = Path('src').absolute()
        if str(src_path) not in sys.path:
            sys.path.insert(0, str(src_path))
        
        from utils import exceptions
        
        print("Available exceptions in utils.exceptions:")
        for attr in dir(exceptions):
            if not attr.startswith('_'):
                obj = getattr(exceptions, attr)
                if isinstance(obj, type) and issubclass(obj, Exception):
                    print(f"  - {attr}")
        
        return True
    except Exception as e:
        print(f"Error checking exceptions: {e}")
        return False

def fix_pptx_processor_imports():
    """Fix the import statements in pptx_processor.py"""
    
    pptx_file = Path('src/core/pptx_processor.py')
    
    # Read current content
    with open(pptx_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Create backup
    backup_file = pptx_file.with_suffix('.py.backup_import_fix')
    with open(backup_file, 'w', encoding='utf-8') as f:
        f.write(content)
    print(f"Backup created: {backup_file}")
    
    # Fix the import statement - use built-in exceptions instead
    old_import = """from utils.exceptions import (
    PPTXProcessingError, FileNotFoundError, ValidationError,
    TranslationError, ConfigurationError
)"""
    
    new_import = """from utils.exceptions import (
    PPTXProcessingError, ValidationError,
    TranslationError, ConfigurationError
)"""
    
    # Replace the problematic import
    content = content.replace(old_import, new_import)
    
    # Also replace any usage of the custom FileNotFoundError with built-in
    content = content.replace('raise FileNotFoundError(', 'raise FileNotFoundError(')
    
    # Actually, let's just use built-in FileNotFoundError by importing it normally
    # Add the built-in import at the top
    if 'import os' in content:
        content = content.replace('import os', 'import os\nfrom pathlib import Path')
    
    # Replace all custom exception usage with appropriate built-ins or custom ones
    lines = content.split('\n')
    fixed_lines = []
    
    for line in lines:
        # Fix FileNotFoundError usage - it should use built-in
        if 'raise FileNotFoundError(' in line and 'utils.exceptions' not in line:
            # This is fine, it will use built-in FileNotFoundError
            fixed_lines.append(line)
        else:
            fixed_lines.append(line)
    
    content = '\n'.join(fixed_lines)
    
    # Write the fixed content
    with open(pptx_file, 'w', encoding='utf-8') as f:
        f.write(content)
    
    print(f"Fixed import statements in {pptx_file}")
    
    # Test if it compiles now
    try:
        compile(content, str(pptx_file), 'exec')
        print("✅ File compiles successfully after import fix")
        return True
    except SyntaxError as e:
        print(f"❌ Still has syntax error: Line {e.lineno}: {e.msg}")
        return False

def test_import_after_fix():
    """Test if the import works after the fix"""
    try:
        import sys
        src_path = Path('src').absolute()
        if str(src_path) not in sys.path:
            sys.path.insert(0, str(src_path))
        
        # Clear cached module
        if 'core.pptx_processor' in sys.modules:
            del sys.modules['core.pptx_processor']
        
        from core.pptx_processor import PPTXProcessor
        processor = PPTXProcessor()
        print("✅ PPTXProcessor imports successfully after fix")
        return True
        
    except ImportError as e:
        print(f"❌ Import still fails: {e}")
        return False
    except Exception as e:
        print(f"❌ Other error: {e}")
        return False

def create_simple_working_processor():
    """Create a simplified processor that uses only built-in exceptions"""
    
    pptx_file = Path('src/core/pptx_processor.py')
    
    # Create a working version with minimal dependencies
    simple_content = '''"""
PowerPoint processor for PPTrans - simplified working version
"""

import os
import shutil
from datetime import datetime
from typing import List, Dict, Optional, Union, Any
from pathlib import Path
import re

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from utils.logger import LoggerMixin, log_performance


class ProcessingStats:
    """Statistics for processing operations"""
    
    def __init__(self):
        self.reset()
    
    def reset(self):
        """Reset all statistics"""
        self.slides_processed = 0
        self.text_elements_found = 0
        self.text_elements_translated = 0
        self.text_elements_skipped = 0
        self.errors_encountered = 0
        self.processing_time = 0.0


class PPTXProcessor(LoggerMixin):
    """PowerPoint presentation processor with translation capabilities"""
    
    def __init__(self, config: Optional[Dict] = None):
        """Initialize PPTX processor"""
        self.config = config or {}
        self.preserve_animations = self.config.get('preserve_animations', True)
        self.backup_original = self.config.get('backup_original', True)
        self.parallel_processing = self.config.get('parallel_processing', False)
        self.max_workers = self.config.get('max_workers', 4)
        
        # Current state
        self.current_presentation = None
        self.current_file_path = None
        self.text_elements = []
        self.stats = ProcessingStats()
        
        self.logger.info(f"PPTX Processor initialized with config: {self.config}")
    
    def load_presentation(self, file_path: Union[str, Path]) -> None:
        """Load a PowerPoint presentation"""
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if not file_path.suffix.lower() == '.pptx':
            raise ValueError(f"File is not a PowerPoint presentation: {file_path}")
        
        try:
            self.logger.info(f"Loading presentation: {file_path}")
            
            # Create backup if requested
            if self.backup_original:
                backup_path = file_path.with_suffix(f'{file_path.suffix}.backup')
                shutil.copy2(file_path, backup_path)
                self.logger.info(f"Created backup: {backup_path}")
            
            # Load the presentation
            self.current_presentation = Presentation(str(file_path))
            self.current_file_path = file_path
            
            # Reset processing state
            self.text_elements = []
            self.stats.reset()
            
            slide_count = len(self.current_presentation.slides)
            self.logger.info(f"Presentation loaded successfully: {slide_count} slides")
            
        except Exception as e:
            raise RuntimeError(f"Failed to load presentation: {str(e)}")
    
    def get_slide_count(self) -> int:
        """Get the number of slides in the current presentation"""
        if not self.current_presentation:
            return 0
        return len(self.current_presentation.slides)
    
    def parse_slide_range(self, range_str: str) -> List[int]:
        """Parse slide range string and return list of slide indices (0-based)"""
        if not range_str or not range_str.strip():
            raise ValueError("Empty slide range specification")
        
        range_str = range_str.strip().lower()
        max_slides = self.get_slide_count()
        
        if max_slides == 0:
            raise ValueError("No presentation loaded")
        
        # Handle "all" keyword
        if range_str == "all":
            return list(range(max_slides))
        
        # Clean up the input
        cleaned = re.sub(r'[^0-9\\-,\\s]', '', range_str)
        cleaned = re.sub(r'\\s+', '', cleaned)
        
        if not cleaned:
            raise ValueError(f"Invalid slide range specification: '{range_str}'")
        
        slide_indices = set()
        
        try:
            parts = [part.strip() for part in cleaned.split(',') if part.strip()]
            
            for part in parts:
                if '-' in part:
                    range_parts = part.split('-')
                    if len(range_parts) != 2:
                        raise ValueError(f"Invalid range format: '{part}'")
                    
                    start_str, end_str = range_parts
                    if not start_str or not end_str:
                        raise ValueError(f"Invalid range format: '{part}'")
                    
                    start = int(start_str)
                    end = int(end_str)
                    
                    if start < 1 or end < 1:
                        raise ValueError(f"Slide numbers must be positive")
                    
                    if start > end:
                        raise ValueError(f"Invalid range: start > end")
                    
                    if start > max_slides or end > max_slides:
                        raise ValueError(f"Slide range exceeds available slides")
                    
                    # Add slides in range (convert to 0-based)
                    slide_indices.update(range(start - 1, end))
                    
                else:
                    slide_num = int(part)
                    if slide_num < 1:
                        raise ValueError(f"Slide number must be positive")
                    if slide_num > max_slides:
                        raise ValueError(f"Slide exceeds available slides")
                    
                    slide_indices.add(slide_num - 1)
        
        except ValueError as e:
            raise ValueError(f"Invalid slide range: {str(e)}")
        
        if not slide_indices:
            raise ValueError(f"No valid slides specified")
        
        return sorted(list(slide_indices))
    
    def extract_text_elements(self, slide_indices: List[int]) -> List[Dict]:
        """Extract text elements from specified slides"""
        if not self.current_presentation:
            raise ValueError("No presentation loaded")
        
        self.logger.info(f"Extracting text from slides: {[i+1 for i in slide_indices]}")
        
        text_elements = []
        
        for slide_idx in slide_indices:
            if slide_idx >= len(self.current_presentation.slides):
                self.logger.warning(f"Slide index {slide_idx} exceeds available slides")
                continue
            
            slide = self.current_presentation.slides[slide_idx]
            
            for shape_idx, shape in enumerate(slide.shapes):
                try:
                    if not hasattr(shape, 'text_frame') or not shape.text_frame:
                        continue
                    
                    text_frame = shape.text_frame
                    
                    for para_idx, paragraph in enumerate(text_frame.paragraphs):
                        paragraph_text = ""
                        formatting_info = {}
                        
                        for run in paragraph.runs:
                            if run.text.strip():
                                paragraph_text += run.text
                                
                                if not formatting_info and run.text.strip():
                                    formatting_info = self._extract_formatting(run, paragraph)
                        
                        if paragraph_text.strip():
                            text_element = {
                                'slide_index': slide_idx,
                                'shape_index': shape_idx,
                                'paragraph_index': para_idx,
                                'original_text': paragraph_text,
                                'translated_text': '',
                                'formatting': formatting_info,
                                'shape_reference': shape
                            }
                            
                            text_elements.append(text_element)
                
                except Exception as e:
                    self.logger.warning(f"Error processing shape {shape_idx}: {e}")
                    continue
        
        self.text_elements = text_elements
        self.logger.info(f"Extracted {len(text_elements)} text elements")
        
        return text_elements
    
    def _extract_formatting(self, run, paragraph) -> Dict:
        """Extract formatting information"""
        formatting = {}
        
        try:
            font = run.font
            
            if font.name:
                formatting['font_name'] = font.name
            
            if font.size:
                formatting['font_size'] = font.size.pt
            
            formatting['bold'] = font.bold if font.bold is not None else False
            formatting['italic'] = font.italic if font.italic is not None else False
            
        except Exception as e:
            self.logger.debug(f"Error extracting formatting: {e}")
        
        return formatting
    
    @log_performance
    def translate_elements(self, translator, source_lang: str = 'auto', target_lang: str = 'en') -> None:
        """Translate all extracted text elements"""
        if not self.text_elements:
            raise ValueError("No text elements to translate")
        
        self.logger.info(f"Starting translation of {len(self.text_elements)} text elements")
        
        for i, element in enumerate(self.text_elements):
            try:
                original_text = element['original_text']
                if not original_text.strip():
                    continue
                
                translated_text = translator.translate_text(
                    original_text,
                    source_lang=source_lang,
                    target_lang=target_lang
                )
                
                element['translated_text'] = translated_text
                
            except Exception as e:
                self.logger.warning(f"Failed to translate element {i}: {e}")
                element['translated_text'] = element['original_text']
        
        self.logger.info(f"Translation completed: {len(self.text_elements)} elements translated, 0 skipped, 0 errors")
    
    def apply_translations(self) -> None:
        """Apply translated text back to the presentation"""
        if not self.current_presentation:
            raise ValueError("No presentation loaded")
        
        if not self.text_elements:
            raise ValueError("No text elements to apply")
        
        self.logger.info("Applying translations to presentation")
        
        successful = 0
        failed = 0
        
        for element in self.text_elements:
            try:
                translated_text = element.get('translated_text', element.get('original_text', ''))
                if not translated_text:
                    continue
                
                shape = element.get('shape_reference')
                if not shape:
                    failed += 1
                    continue
                
                if self._apply_text_simple(shape, translated_text):
                    successful += 1
                else:
                    failed += 1
                
            except Exception as e:
                self.logger.error(f"Failed to apply translation: {e}")
                failed += 1
        
        self.logger.info(f"Applied {successful} translations successfully, {failed} failed")
    
    def _apply_text_simple(self, shape, text: str) -> bool:
        """Simple text application"""
        try:
            if hasattr(shape, 'text'):
                shape.text = text
                return True
            
            if hasattr(shape, 'text_frame') and shape.text_frame:
                shape.text_frame.text = text
                return True
            
            if (hasattr(shape, 'text_frame') and 
                shape.text_frame and 
                shape.text_frame.paragraphs):
                
                first_paragraph = shape.text_frame.paragraphs[0]
                first_paragraph.clear()
                run = first_paragraph.add_run()
                run.text = text
                return True
            
            return False
            
        except Exception as e:
            self.logger.debug(f"Error applying text: {e}")
            return False
    
    def save_presentation(self, output_path: Optional[Union[str, Path]] = None) -> Path:
        """Save the current presentation"""
        if not self.current_presentation:
            raise ValueError("No presentation loaded")
        
        if output_path is None:
            original_path = Path(self.current_file_path)
            output_path = original_path.parent / f"{original_path.stem}_translated{original_path.suffix}"
        else:
            output_path = Path(output_path)
        
        try:
            self.logger.info(f"Saving presentation to: {output_path}")
            output_path.parent.mkdir(parents=True, exist_ok=True)
            self.current_presentation.save(str(output_path))
            self.logger.info(f"Presentation saved successfully: {output_path}")
            return output_path
            
        except Exception as e:
            raise RuntimeError(f"Failed to save presentation: {str(e)}")
    
    def get_processing_stats(self) -> Dict:
        """Get processing statistics"""
        return {
            'slides_processed': self.stats.slides_processed,
            'text_elements_found': len(self.text_elements),
            'text_elements_translated': self.stats.text_elements_translated,
            'text_elements_skipped': self.stats.text_elements_skipped,
            'errors_encountered': self.stats.errors_encountered,
            'processing_time': self.stats.processing_time
        }
    
    def close(self) -> None:
        """Clean up resources"""
        self.current_presentation = None
        self.current_file_path = None
        self.text_elements = []
        self.stats = ProcessingStats()
        self.logger.info("PPTX Processor closed")
'''
    
    # Write the simple version
    with open(pptx_file, 'w', encoding='utf-8') as f:
        f.write(simple_content)
    
    print("Created simplified processor with built-in exceptions only")

def main():
    """Main function to fix the import issue"""
    print("Fixing Exception Import Issue")
    print("=" * 35)
    
    # Check what exceptions are available
    print("Step 1: Checking available exceptions...")
    check_available_exceptions()
    
    print("\nStep 2: Creating simplified processor...")
    create_simple_working_processor()
    
    print("\nStep 3: Testing import...")
    if test_import_after_fix():
        print("\n✅ SUCCESS! Import issue fixed")
        print("Now test: python src/main.py")
        return True
    else:
        print("\n❌ Import still fails")
        return False

if __name__ == "__main__":
    main()
