"""
PPTX Processor with batch translation support
"""
import os
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple, Callable
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor
from utils.logger import LoggerMixin
from utils.exceptions import ValidationError, PPTXProcessingError


class PPTXProcessor(LoggerMixin):
    """PowerPoint processor optimized for batch translation"""
    
    def __init__(self, advanced_settings: dict):
        """Initialize with advanced settings from config"""
        self.settings = advanced_settings
        self.presentation = None
        self.current_file = None
        self.text_elements = []
        self.processing_stats = {
            'total_slides': 0,
            'processed_slides': 0,
            'text_elements_found': 0,
            'translated_elements': 0,
            'skipped_elements': 0,
            'error_count': 0
        }
        
        self.logger.info("Batch-enabled PPTXProcessor initialized")
        
    def load_presentation(self, file_path: str):
        """Load the PowerPoint presentation"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
            
        try:
            self.current_file = file_path
            self.presentation = Presentation(file_path)
            self.processing_stats['total_slides'] = len(self.presentation.slides)
            self.logger.info(f"Loaded presentation with {self.processing_stats['total_slides']} slides")
        except Exception as e:
            self.logger.error(f"Failed to load presentation: {e}")
            raise PPTXProcessingError(f"Cannot load PowerPoint file: {e}", file_path=file_path)
    
    def get_presentation_info(self) -> Dict[str, Any]:
        """Get presentation information"""
        if not self.presentation:
            return {'total_slides': 0}
        
        return {
            'total_slides': len(self.presentation.slides),
            'file_path': self.current_file
        }
    
    def parse_slide_range(self, slide_range: str) -> List[int]:
        """Parse slide range specification into list of slide indices"""
        if not self.presentation:
            raise PPTXProcessingError("No presentation loaded")
        
        total_slides = len(self.presentation.slides)
        
        if slide_range.lower() == 'all':
            return list(range(total_slides))
        
        indices = []
        parts = slide_range.split(',')
        
        for part in parts:
            part = part.strip()
            if not part:
                continue
            
            if '-' in part:
                try:
                    start_str, end_str = part.split('-', 1)
                    start = int(start_str.strip()) - 1  # Convert to 0-based
                    end = int(end_str.strip()) - 1      # Convert to 0-based
                    
                    # Validate range
                    start = max(0, min(start, total_slides - 1))
                    end = max(0, min(end, total_slides - 1))
                    
                    if start <= end:
                        indices.extend(range(start, end + 1))
                    else:
                        indices.extend(range(end, start + 1))
                        
                except ValueError as e:
                    self.logger.error(f"Invalid range format: {part}")
                    continue
            else:
                try:
                    slide_num = int(part) - 1  # Convert to 0-based
                    if 0 <= slide_num < total_slides:
                        indices.append(slide_num)
                except ValueError:
                    self.logger.error(f"Invalid slide number: {part}")
                    continue
        
        # Remove duplicates and sort
        indices = sorted(list(set(indices)))
        self.logger.debug(f"Parsed slide range '{slide_range}' to indices: {indices}")
        return indices
    
    def _extract_run_formatting(self, run) -> Dict[str, Any]:
        """Extract all formatting information from a text run"""
        formatting = {}
        
        try:
            if hasattr(run, 'font'):
                font = run.font
                
                # Font name
                if hasattr(font, 'name') and font.name:
                    formatting['font_name'] = font.name
                
                # Font size
                if hasattr(font, 'size') and font.size:
                    formatting['font_size'] = font.size.pt
                
                # Bold, italic, underline
                formatting['bold'] = font.bold if font.bold is not None else False
                formatting['italic'] = font.italic if font.italic is not None else False
                formatting['underline'] = font.underline if font.underline is not None else False
                
                # Font color
                try:
                    if hasattr(font, 'color') and font.color and hasattr(font.color, 'rgb'):
                        rgb = font.color.rgb
                        if rgb:
                            formatting['font_color'] = (rgb.r, rgb.g, rgb.b)
                except Exception:
                    pass  # Color extraction can fail
                    
        except Exception as e:
            self.logger.debug(f"Error extracting run formatting: {e}")
        
        return formatting
    
    def _extract_paragraph_formatting(self, paragraph) -> Dict[str, Any]:
        """Extract paragraph-level formatting"""
        formatting = {}
        
        try:
            # Paragraph alignment
            if hasattr(paragraph, 'alignment') and paragraph.alignment is not None:
                alignment_map = {
                    PP_ALIGN.LEFT: 'left',
                    PP_ALIGN.CENTER: 'center', 
                    PP_ALIGN.RIGHT: 'right',
                    PP_ALIGN.JUSTIFY: 'justify'
                }
                formatting['alignment'] = alignment_map.get(paragraph.alignment, 'left')
            
            # Space before/after
            if hasattr(paragraph, 'space_before') and paragraph.space_before:
                formatting['space_before'] = paragraph.space_before.pt
                
            if hasattr(paragraph, 'space_after') and paragraph.space_after:
                formatting['space_after'] = paragraph.space_after.pt
                
        except Exception as e:
            self.logger.debug(f"Error extracting paragraph formatting: {e}")
        
        return formatting
    
    def extract_text_elements(self, slide_range: str = "all") -> List[Dict[str, Any]]:
        """
        Extract text elements from specified slides with smart grouping for batch translation
        """
        if not self.presentation:
            raise PPTXProcessingError("No presentation loaded")
        
        slide_indices = self.parse_slide_range(slide_range)
        self.text_elements = []
        element_id = 0
        
        self.logger.info(f"Extracting text from slides: {[i+1 for i in slide_indices]}")
        
        for slide_idx in slide_indices:
            try:
                slide = self.presentation.slides[slide_idx]
                slide_elements = 0
                
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        text_frame = shape.text_frame
                        
                        # Group by paragraph for better context in batch translation
                        for para_idx, paragraph in enumerate(text_frame.paragraphs):
                            if not paragraph.text.strip():
                                continue
                                
                            paragraph_formatting = self._extract_paragraph_formatting(paragraph)
                            
                            # Collect all runs in this paragraph
                            paragraph_runs = []
                            paragraph_text_parts = []
                            
                            for run_idx, run in enumerate(paragraph.runs):
                                if run.text:  # Include even empty runs to preserve structure
                                    run_formatting = self._extract_run_formatting(run)
                                    
                                    run_data = {
                                        'run_index': run_idx,
                                        'text': run.text,
                                        'formatting': run_formatting,
                                        'run_ref': run
                                    }
                                    paragraph_runs.append(run_data)
                                    paragraph_text_parts.append(run.text)
                            
                            if paragraph_runs:  # Only if we have runs with text
                                # Create a paragraph-level element for batch translation
                                element = {
                                    'id': element_id,
                                    'slide_index': slide_idx,
                                    'shape_index': shape_idx,
                                    'paragraph_index': para_idx,
                                    'type': 'paragraph',
                                    'original_text': ''.join(paragraph_text_parts),
                                    'translated_text': None,
                                    'paragraph_formatting': paragraph_formatting,
                                    'runs': paragraph_runs,
                                    'shape_ref': shape,
                                    'paragraph_ref': paragraph
                                }
                                
                                self.text_elements.append(element)
                                element_id += 1
                                slide_elements += 1
                
                self.processing_stats['processed_slides'] += 1
                self.logger.debug(f"Slide {slide_idx + 1}: Found {slide_elements} paragraph elements")
                
            except Exception as e:
                self.logger.error(f"Error processing slide {slide_idx + 1}: {e}")
                self.processing_stats['error_count'] += 1
        
        self.processing_stats['text_elements_found'] = len(self.text_elements)
        self.logger.info(f"Extracted {len(self.text_elements)} paragraph elements from {len(slide_indices)} slides")
        
        return self.text_elements
    
    def translate_text_elements(self, translate_callback: Callable[[str], str]) -> None:
        """
        Translate extracted text elements using batch processing
        
        Args:
            translate_callback: Function that handles translation - now expects to handle batches
        """
        if not self.text_elements:
            raise ValidationError("No text elements to translate")
        
        # Check if the translator supports batch translation
        from core.translator import PPTransTranslator
        translator = PPTransTranslator(self.settings.get('translation', {}))
        
        # Collect all text to translate
        texts_to_translate = []
        element_map = {}  # Map translation results back to elements
        
        for i, element in enumerate(self.text_elements):
            original_text = element['original_text']
            if original_text and original_text.strip():
                texts_to_translate.append(original_text)
                element_map[len(texts_to_translate) - 1] = i
            else:
                element['translated_text'] = original_text  # Keep empty/whitespace as is
        
        if not texts_to_translate:
            self.logger.info("No text found that needs translation")
            return
        
        self.logger.info(f"Starting batch translation of {len(texts_to_translate)} text elements")
        
        try:
            # Use batch translation
            translated_texts = translator.translate_text_batch(
                texts_to_translate, 
                source_lang='de', 
                target_lang='en'
            )
            
            # Map results back to elements
            translated_count = 0
            for batch_idx, translated_text in enumerate(translated_texts):
                if batch_idx in element_map:
                    element_idx = element_map[batch_idx]
                    element = self.text_elements[element_idx]
                    
                    if translated_text and translated_text != element['original_text']:
                        element['translated_text'] = translated_text
                        translated_count += 1
                        self.logger.debug(f"Translated paragraph {element['id']}: '{element['original_text'][:50]}...' -> '{translated_text[:50]}...'")
                    else:
                        element['translated_text'] = element['original_text']
            
            self.processing_stats['translated_elements'] = translated_count
            self.processing_stats['skipped_elements'] = len(self.text_elements) - translated_count
            
            self.logger.info(f"Batch translation completed: {translated_count}/{len(self.text_elements)} elements translated")
            
        except Exception as e:
            self.logger.error(f"Batch translation failed: {e}")
            # Fall back to original text for all elements
            for element in self.text_elements:
                element['translated_text'] = element['original_text']
            self.processing_stats['error_count'] += 1
    
    def _distribute_translation_to_runs(self, element: Dict[str, Any]) -> None:
        """
        Distribute translated paragraph text back to individual runs
        This is the tricky part - we need to map the translated text back to the original run structure
        """
        original_text = element['original_text']
        translated_text = element['translated_text']
        runs = element['runs']
        
        if not translated_text or translated_text == original_text:
            # No translation occurred, keep original run texts
            return
        
        # Simple approach: if translation is close in length, distribute proportionally
        # More sophisticated approach would use word alignment
        
        original_length = len(original_text)
        translated_length = len(translated_text)
        
        if original_length == 0:
            return
        
        # Calculate proportional distribution
        distributed_texts = []
        current_pos = 0
        
        for run_data in runs:
            run_text = run_data['text']
            run_length = len(run_text)
            
            if run_length == 0:
                distributed_texts.append('')
                continue
            
            # Calculate proportion of translated text for this run
            proportion = run_length / original_length
            target_length = int(translated_length * proportion)
            
            # Extract portion of translated text
            if current_pos + target_length <= translated_length:
                run_translated = translated_text[current_pos:current_pos + target_length]
                current_pos += target_length
            else:
                # Take remaining text for last run
                run_translated = translated_text[current_pos:]
                current_pos = translated_length
            
            # Handle word boundaries to avoid splitting words
            if run_translated and current_pos < translated_length:
                # Try to end at word boundary
                last_space = run_translated.rfind(' ')
                if last_space > len(run_translated) * 0.5:  # Only if we're not cutting too much
                    run_translated = run_translated[:last_space + 1]
                    current_pos = current_pos - (len(run_translated) - last_space - 1)
            
            distributed_texts.append(run_translated)
        
        # If we have remaining text, append to last run
        if current_pos < translated_length:
            remaining = translated_text[current_pos:]
            if distributed_texts:
                distributed_texts[-1] += remaining
        
        # Update run data with distributed translations
        for i, run_data in enumerate(runs):
            if i < len(distributed_texts):
                run_data['translated_text'] = distributed_texts[i]
            else:
                run_data['translated_text'] = run_data['text']  # Fallback to original
    
    def _apply_run_formatting(self, run, formatting: Dict[str, Any]) -> None:
        """Apply formatting to a text run"""
        try:
            if hasattr(run, 'font'):
                font = run.font
                
                # Apply font name
                if 'font_name' in formatting and formatting['font_name']:
                    font.name = formatting['font_name']
                
                # Apply font size
                if 'font_size' in formatting and formatting['font_size']:
                    font.size = Pt(formatting['font_size'])
                
                # Apply bold, italic, underline
                if 'bold' in formatting:
                    font.bold = formatting['bold']
                if 'italic' in formatting:
                    font.italic = formatting['italic']
                if 'underline' in formatting:
                    font.underline = formatting['underline']
                
                # Apply font color
                if 'font_color' in formatting and formatting['font_color']:
                    r, g, b = formatting['font_color']
                    font.color.rgb = RGBColor(r, g, b)
                    
        except Exception as e:
            self.logger.debug(f"Error applying run formatting: {e}")
    
    def _apply_paragraph_formatting(self, paragraph, formatting: Dict[str, Any]) -> None:
        """Apply formatting to a paragraph"""
        try:
            # Apply alignment
            if 'alignment' in formatting:
                alignment_map = {
                    'left': PP_ALIGN.LEFT,
                    'center': PP_ALIGN.CENTER,
                    'right': PP_ALIGN.RIGHT,
                    'justify': PP_ALIGN.JUSTIFY
                }
                if formatting['alignment'] in alignment_map:
                    paragraph.alignment = alignment_map[formatting['alignment']]
            
            # Apply spacing
            if 'space_before' in formatting:
                paragraph.space_before = Pt(formatting['space_before'])
            if 'space_after' in formatting:
                paragraph.space_after = Pt(formatting['space_after'])
                
        except Exception as e:
            self.logger.debug(f"Error applying paragraph formatting: {e}")
    
    def apply_translations(self) -> None:
        """Apply translated text to the presentation with preserved formatting"""
        if not self.text_elements:
            raise ValidationError("No text elements to apply")
        
        self.logger.info("Applying batch translations to presentation")
        
        applied_count = 0
        error_count = 0
        
        for element in self.text_elements:
            try:
                if element['translated_text'] is None:
                    continue
                
                # Distribute translation back to runs
                self._distribute_translation_to_runs(element)
                
                # Apply translated text to each run
                for run_data in element['runs']:
                    run = run_data['run_ref']
                    translated_text = run_data.get('translated_text', run_data['text'])
                    
                    # Apply the translated text
                    run.text = translated_text
                    
                    # Reapply formatting
                    self._apply_run_formatting(run, run_data['formatting'])
                
                # Apply paragraph formatting
                self._apply_paragraph_formatting(element['paragraph_ref'], element['paragraph_formatting'])
                
                applied_count += 1
                self.logger.debug(f"Applied translation to paragraph {element['id']}")
                
            except Exception as e:
                self.logger.error(f"Error applying translation to element {element.get('id', 'unknown')}: {e}")
                error_count += 1
        
        self.logger.info(f"Applied {applied_count} translations, {error_count} errors")
        self.processing_stats['error_count'] += error_count
    
    def save_presentation(self) -> str:
        """Save the modified presentation"""
        if not self.presentation or not self.current_file:
            raise PPTXProcessingError("No presentation loaded to save")
            
        try:
            # Create output filename
            input_path = Path(self.current_file)
            output_path = input_path.parent / f"{input_path.stem}_translated{input_path.suffix}"
            
            self.presentation.save(str(output_path))
            self.logger.info(f"Presentation saved to: {output_path}")
            return str(output_path)
        except Exception as e:
            self.logger.error(f"Failed to save presentation: {e}")
            raise PPTXProcessingError(f"Cannot save presentation: {e}", file_path=str(output_path))
    
    def get_processing_stats(self) -> Dict[str, Any]:
        """Get processing statistics"""
        return self.processing_stats.copy()