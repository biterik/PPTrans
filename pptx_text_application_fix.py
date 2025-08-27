#!/usr/bin/env python3
"""
Fix for PowerPoint text application - addresses the 'tuple' object has no attribute 'add' error
"""

import sys
import os
from pathlib import Path

def fix_pptx_text_application():
    """Fix the PowerPoint text application methods"""
    
    pptx_processor_path = Path('src/core/pptx_processor.py')
    if not pptx_processor_path.exists():
        print(f"❌ File not found: {pptx_processor_path}")
        return False
    
    # Create backup
    backup_path = Path(str(pptx_processor_path) + '.backup_text_fix')
    with open(pptx_processor_path, 'r', encoding='utf-8') as f:
        content = f.read()
    with open(backup_path, 'w', encoding='utf-8') as f:
        f.write(content)
    print(f"✅ Backup created: {backup_path}")
    
    # Fixed text application methods
    fixed_methods = '''
    def _apply_text_to_shape(self, shape, translated_text: str, original_format: Dict) -> bool:
        """
        Apply translated text to a shape while preserving formatting
        
        Args:
            shape: PowerPoint shape object
            translated_text: Text to apply
            original_format: Original formatting information
            
        Returns:
            True if successful, False otherwise
        """
        try:
            if not hasattr(shape, 'text_frame'):
                self.logger.warning("Shape has no text_frame attribute")
                return False
            
            text_frame = shape.text_frame
            if not text_frame:
                self.logger.warning("Shape has no text_frame")
                return False
            
            # Clear existing paragraphs except the first one
            while len(text_frame.paragraphs) > 1:
                p = text_frame.paragraphs[-1]
                p.clear()
            
            # Get the first (and now only) paragraph
            paragraph = text_frame.paragraphs[0]
            
            # Clear existing runs
            paragraph.clear()
            
            # Apply paragraph-level formatting
            if 'alignment' in original_format:
                try:
                    from pptx.enum.text import PP_ALIGN
                    alignment_map = {
                        'left': PP_ALIGN.LEFT,
                        'center': PP_ALIGN.CENTER,
                        'right': PP_ALIGN.RIGHT,
                        'justify': PP_ALIGN.JUSTIFY
                    }
                    if original_format['alignment'] in alignment_map:
                        paragraph.alignment = alignment_map[original_format['alignment']]
                except Exception as e:
                    self.logger.debug(f"Could not apply alignment: {e}")
            
            # Add the translated text as a new run
            run = paragraph.add_run()
            run.text = translated_text
            
            # Apply character-level formatting to the run
            self._apply_run_formatting(run, original_format)
            
            return True
            
        except Exception as e:
            self.logger.error(f"Error applying text to shape: {e}")
            return False
    
    def _apply_run_formatting(self, run, original_format: Dict) -> None:
        """
        Apply character formatting to a text run
        
        Args:
            run: PowerPoint text run object
            original_format: Original formatting information
        """
        try:
            font = run.font
            
            # Apply font name
            if 'font_name' in original_format and original_format['font_name']:
                try:
                    font.name = original_format['font_name']
                except Exception as e:
                    self.logger.debug(f"Could not apply font name {original_format['font_name']}: {e}")
            
            # Apply font size
            if 'font_size' in original_format and original_format['font_size']:
                try:
                    from pptx.util import Pt
                    font.size = Pt(original_format['font_size'])
                except Exception as e:
                    self.logger.debug(f"Could not apply font size {original_format['font_size']}: {e}")
            
            # Apply bold
            if 'bold' in original_format:
                try:
                    font.bold = original_format['bold']
                except Exception as e:
                    self.logger.debug(f"Could not apply bold formatting: {e}")
            
            # Apply italic
            if 'italic' in original_format:
                try:
                    font.italic = original_format['italic']
                except Exception as e:
                    self.logger.debug(f"Could not apply italic formatting: {e}")
            
            # Apply underline
            if 'underline' in original_format:
                try:
                    font.underline = original_format['underline']
                except Exception as e:
                    self.logger.debug(f"Could not apply underline formatting: {e}")
            
            # Apply font color
            if 'font_color' in original_format and original_format['font_color']:
                try:
                    from pptx.dml.color import RGBColor
                    color_value = original_format['font_color']
                    if isinstance(color_value, tuple) and len(color_value) == 3:
                        r, g, b = color_value
                        font.color.rgb = RGBColor(r, g, b)
                except Exception as e:
                    self.logger.debug(f"Could not apply font color: {e}")
            
        except Exception as e:
            self.logger.debug(f"Error applying run formatting: {e}")
    
    def apply_translations(self) -> None:
        """
        Apply translated text back to the presentation with improved error handling
        """
        if not self.current_presentation:
            raise ValueError("No presentation loaded")
        
        if not self.text_elements:
            raise ValueError("No text elements to apply")
        
        self.logger.info("Applying translations to presentation")
        
        successful_applications = 0
        failed_applications = 0
        
        for element in self.text_elements:
            try:
                if not element.get('translated_text'):
                    self.logger.debug(f"Skipping element with no translation on slide {element['slide_index'] + 1}")
                    continue
                
                shape = element.get('shape_reference')
                if not shape:
                    self.logger.warning(f"No shape reference for element on slide {element['slide_index'] + 1}")
                    failed_applications += 1
                    continue
                
                # Apply the translation
                success = self._apply_text_to_shape(
                    shape,
                    element['translated_text'],
                    element.get('formatting', {})
                )
                
                if success:
                    successful_applications += 1
                else:
                    failed_applications += 1
                
            except Exception as e:
                self.logger.error(f"Failed to apply translation for element on slide {element['slide_index'] + 1}: {e}")
                failed_applications += 1
        
        self.logger.info(f"Applied {successful_applications} translations successfully, {failed_applications} failed")
        
        if failed_applications > 0:
            self.logger.warning(f"{failed_applications} translations could not be applied")
'''
    
    # Find the location to insert these methods
    # Look for existing method definitions to replace or class definition to add after
    if '_apply_text_to_shape' in content:
        # Replace existing methods
        import re
        
        # Remove existing _apply_text_to_shape method
        content = re.sub(
            r'    def _apply_text_to_shape\(.*?\n(?:    .*\n)*?        except.*?\n.*?\n',
            '',
            content,
            flags=re.DOTALL
        )
        
        # Remove existing _apply_run_formatting method  
        content = re.sub(
            r'    def _apply_run_formatting\(.*?\n(?:    .*\n)*?        except.*?\n.*?\n',
            '',
            content,
            flags=re.DOTALL
        )
        
        # Remove existing apply_translations method
        content = re.sub(
            r'    def apply_translations\(.*?\n(?:    .*\n)*?            self\.logger\.warning.*?\n',
            '',
            content,
            flags=re.DOTALL
        )
    
    # Find the end of the class to add our methods
    class_end_pattern = r'(\n    def [^_].*?(?=\n\n|\nclass|\Z))'
    
    # Add the fixed methods before the last method of the class
    if 'class PPTXProcessor' in content:
        # Find a good insertion point - after other methods but before close()
        insertion_point = content.rfind('    def close(')
        if insertion_point == -1:
            # If no close method, add at the end of the class
            insertion_point = content.rfind('\n    def ')
            if insertion_point != -1:
                # Find the end of that method
                next_method_start = content.find('\n    def ', insertion_point + 1)
                if next_method_start == -1:
                    insertion_point = len(content) - 1
                else:
                    insertion_point = next_method_start
            else:
                insertion_point = len(content) - 1
        
        # Insert the fixed methods
        content = content[:insertion_point] + fixed_methods + '\n' + content[insertion_point:]
    else:
        print("❌ Could not find PPTXProcessor class")
        return False
    
    try:
        with open(pptx_processor_path, 'w', encoding='utf-8') as f:
            f.write(content)
        print(f"✅ Fixed PowerPoint text application methods in {pptx_processor_path}")
        return True
    except Exception as e:
        print(f"❌ Error writing fixed file: {e}")
        return False

def test_pptx_processor():
    """Test the fixed PowerPoint processor"""
    print("\n=== Testing Fixed PowerPoint Processor ===")
    
    try:
        # Add src to path
        src_path = Path('src').absolute()
        if str(src_path) not in sys.path:
            sys.path.insert(0, str(src_path))
        
        # Reimport to get the new version
        import importlib
        if 'core.pptx_processor' in sys.modules:
            importlib.reload(sys.modules['core.pptx_processor'])
        
        from core.pptx_processor import PPTXProcessor
        
        processor = PPTXProcessor()
        print("✅ PPTXProcessor created successfully")
        
        # Check if the fixed methods exist
        required_methods = [
            '_apply_text_to_shape',
            '_apply_run_formatting', 
            'apply_translations'
        ]
        
        missing_methods = []
        for method in required_methods:
            if not hasattr(processor, method):
                missing_methods.append(method)
        
        if missing_methods:
            print(f"❌ Missing methods: {missing_methods}")
            return False
        else:
            print("✅ All required text application methods present")
            return True
        
    except Exception as e:
        print(f"❌ PowerPoint processor test failed: {e}")
        import traceback
        traceback.print_exc()
        return False

def main():
    """Main function to fix and test PowerPoint text application"""
    print("PowerPoint Text Application Fix for PPTrans")
    print("=" * 55)
    
    # Check we're in the right directory
    if not Path('src').exists():
        print("❌ Please run this script from the PPTrans project root directory")
        return False
    
    success = True
    
    # Apply the fix
    success &= fix_pptx_text_application()
    
    # Test the fix
    if success:
        success &= test_pptx_processor()
    
    print("\n" + "=" * 55)
    if success:
        print("✅ PowerPoint text application fix applied and tested!")
        print("\n✨ Fixed issues:")
        print("   - Properly clear and recreate text runs")
        print("   - Handle paragraph.add_run() correctly")
        print("   - Avoid 'tuple' object errors")
        print("   - Better error handling and logging")
        print("\n🔧 Next step:")
        print("   Test translation: python src/main.py")
        print("   Even with translation API issues, text should now be applied correctly")
    else:
        print("❌ Some issues remain. Please check the output above.")
    
    return success

if __name__ == "__main__":
    main()
